from pptx import Presentation
import io

def parse_text_to_ppt(ai_text):
    """마크다운 텍스트를 분석하여 PPT 버퍼로 반환"""
    prs = Presentation()
    lines = ai_text.split('\n')
    
    current_title = "OTGALNON AI 분석 결과"
    current_bullets = []
    
    slide_layout_title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout_title)
    slide.shapes.title.text = "OTGALNON 실시간 발표자료"
    slide.placeholders[1].text = "AI가 실시간으로 추출한 데이터 기반 보고서"

    slide_layout_content = prs.slide_layouts[1]

    def add_slide_to_presentation(title, bullets):
        if bullets or title != "OTGALNON AI 분석 결과":
            s = prs.slides.add_slide(slide_layout_content)
            s.shapes.title.text = title.replace('#', '').replace('*', '').strip()
            tf = s.placeholders[1].text_frame
            if bullets:
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    tf.add_paragraph().text = bullet

    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        if line.startswith('#'):
            add_slide_to_presentation(current_title, current_bullets)
            current_title = line
            current_bullets = []
        elif line.startswith(('*', '-', '•')) or (line[0].isdigit() if len(line) > 0 else False):
            clean_bullet = line.lstrip('*-•0123456789. ')
            if clean_bullet:
                current_bullets.append(clean_bullet)
        else:
            if len(current_bullets) < 5:
                current_bullets.append(line)
                
    add_slide_to_presentation(current_title, current_bullets)

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer