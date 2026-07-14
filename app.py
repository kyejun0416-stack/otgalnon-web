import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import io
import json
import os

# ==========================================
# 1. 페이지 및 테마 설정
# ==========================================
st.set_page_config(page_title="OTGALNON", page_icon="logo.png", layout="wide")

st.markdown("""
    <style>
    /* 아바타 숨김 및 여백 최적화 */
    [data-testid="stChatMessageAvatar"],
    .stChatMessageAvatar,
    div[data-testid="chatAvatarIcon-user"],
    div[data-testid="chatAvatarIcon-assistant"] {
        display: none !important;
        width: 0px !important;
        height: 0px !important;
        margin: 0px !important;
        padding: 0px !important;
    }
    [data-testid="stChatMessage"] {
        padding: 1rem 0 !important;
        gap: 0rem !important;
    }
    [data-testid="stChatMessageContent"] {
        margin-left: 0px !important;
        padding-left: 0.5rem !important;
    }
    [data-testid="stSidebar"] img {
        margin-bottom: 2rem;
        border-radius: 8px;
    }
    code {
        color: #b39ddb !important;
    }
    </style>
    """, unsafe_allow_html=True)

# ==========================================
# 2. 대화 기록 자동 저장 및 복구 로직 (File I/O)
# ==========================================
HISTORY_FILE = "chat_history.json"

def load_history():
    """파일에서 대화 기록을 불러옵니다."""
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return []
    return []

def save_history(messages):
    """대화 기록을 파일에 즉시 저장합니다."""
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(messages, f, ensure_ascii=False, indent=4)

# ==========================================
# 3. AI 답변 기반 동적 PPT 생성 로직
# ==========================================
def parse_text_to_ppt(ai_text):
    prs = Presentation()
    lines = ai_text.split('\n')
    
    current_title = "OTGALNON AI 분석 결과"
    current_bullets = []
    
    slide_layout_title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout_title)
    slide.shapes.title.text = "OTGALNON 실시간 발표자료"
    slide.placeholders[1].text = "AI가 실시간으로 추출한 데이터 기반 보고서"

    slide_layout_content = prs.slide_layouts[1]

    def add_slide_to_presentation(title, bullets):
        if bullets or title != "OTGALNON AI 분석 결과":
            s = prs.slides.add_slide(slide_layout_content)
            s.shapes.title.text = title.replace('#', '').replace('*', '').strip()
            tf = s.placeholders[1].text_frame
            if bullets:
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    tf.add_paragraph().text = bullet

    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        if line.startswith('#'):
            add_slide_to_presentation(current_title, current_bullets)
            current_title = line
            current_bullets = []
        elif line.startswith(('*', '-', '•')) or (line[0].isdigit() if len(line) > 0 else False):
            clean_bullet = line.lstrip('*-•0123456789. ')
            if clean_bullet:
                current_bullets.append(clean_bullet)
        else:
            if len(current_bullets) < 5:
                current_bullets.append(line)
                
    add_slide_to_presentation(current_title, current_bullets)

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# ==========================================
# 4. 시스템 엔진 초기화
# ==========================================
def initialize_otgalnon():
    api_key = st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        return None, "API Key Missing"
    
    genai.configure(api_key=api_key)
    try:
        model_list = [m.name for m in genai.list_models()]
        target = next((m for m in model_list if "gemini-3-flash" in m), "models/gemini-1.5-flash")
        
        model = genai.GenerativeModel(
            model_name=target,
            system_instruction=(
                "당신은 OTGALNON의 최고 분석관입니다. "
                "이모티콘 사용을 엄격히 금지하며, 제1원리 추론에 기반해 논리적으로 답변하십시오. "
                "사용자가 발표 자료나 PPT 형태의 출력을 요구하면, 각 슬라이드의 제목은 '###' 또는 '##'으로 구분하고, "
                "내용은 글머리 기호(* 또는 -)를 사용하여 요약식으로 가독성 있게 작성하십시오."
            )
        )
        return model, target
    except Exception as e:
        return None, str(e)

engine, active_id = initialize_otgalnon()

# ==========================================
# 5. 사이드바 컨트롤 패널
# ==========================================
with st.sidebar:
    try:
        st.image("logo.png", use_column_width=True)
    except:
        st.markdown("<h2 style='color: #6d5dfc; text-align: center; margin-bottom: 2rem;'>OTGALNON</h2>", unsafe_allow_html=True)
    
    st.markdown("### SYSTEM STATUS")
    if engine:
        st.success("ONLINE")
    else:
        st.error("OFFLINE")
    
    # [수정됨] 워크스페이스 리셋 시 파일도 함께 삭제하여 완벽하게 비움
    if st.button("RESET WORKSPACE", use_container_width=True):
        st.session_state.messages = []
        if os.path.exists(HISTORY_FILE):
            os.remove(HISTORY_FILE)
        st.rerun()

    st.markdown("---")
    st.markdown("### 🚀 OTGALNON TOOLS")
    
    last_assistant_response = None
    if "messages" in st.session_state:
        for msg in reversed(st.session_state.messages):
            if msg["role"] == "assistant":
                last_assistant_response = msg["content"]
                break

    if last_assistant_response:
        try:
            ppt_file = parse_text_to_ppt(last_assistant_response)
            st.download_button(
                label="📥 마지막 답변 PPT로 다운로드",
                data=ppt_file,
                file_name="OTGALNON_AI_발표자료.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        except Exception as e:
            st.caption(f"PPT 파일 변환 대기 중...")
    else:
        st.info("오트가논에게 명령을 내려서 답변을 받아보세요. 그 답변으로 PPT를 만들어 드립니다.")

# ==========================================
# 6. 메인 챗봇 인터페이스
# ==========================================
# [수정됨] 시작할 때 파일에 저장된 기록이 있으면 불러옴
if "messages" not in st.session_state:
    st.session_state.messages = load_history()

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

if prompt := st.chat_input("명령을 입력하십시오..."):
    # 사용자 메시지 저장
    st.session_state.messages.append({"role": "user", "content": prompt})
    save_history(st.session_state.messages)
    
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        if engine:
            with st.spinner("분석 중..."):
                try:
                    response = engine.generate_content(prompt)
                    answer = response.text
                    st.markdown(answer)
                    
                    # 어시스턴트(AI) 메시지 저장
                    st.session_state.messages.append({"role": "assistant", "content": answer})
                    save_history(st.session_state.messages)
                    
                    st.rerun() # 버튼 상태 갱신
                except Exception as e:
                    st.error(f"시스템 오류: {str(e)}")
        else:
            st.warning("엔진 연결을 확인하십시오.")